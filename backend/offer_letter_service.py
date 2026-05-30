"""
Offer Letter Service — Cross-Platform

Core logic for:
1. Filling PPTX templates with employee data
2. Converting PPTX to PDF via LibreOffice headless (replaces win32com)
3. Sending email via Gmail API with PDF attachment
"""

import os
import subprocess
import platform
import shutil
import base64
import mimetypes
import time
from email.message import EmailMessage

from pptx import Presentation
from gmail_auth import get_gmail_service

# =====================================================
# CONFIG
# =====================================================

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATES_DIR = os.path.join(BASE_DIR, 'templates')
OUTPUT_DIR = os.path.join(BASE_DIR, 'output')
FONT_NAME = "Poppins"

# Ensure output directory exists
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Email config — loaded from environment or defaults
SENDER_EMAIL = os.environ.get('SENDER_EMAIL', 'sairamjoshi.cs@gmail.com')
CC_EMAIL = os.environ.get('CC_EMAIL', 'support@collegesimplified.in')


# =====================================================
# FIND LIBREOFFICE BINARY
# =====================================================

def _find_libreoffice():
    """Find the LibreOffice binary path based on the OS."""
    system = platform.system()

    if system == 'Windows':
        # Common Windows install paths
        candidates = [
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
        ]
        for path in candidates:
            if os.path.exists(path):
                return path
        # Try PATH
        result = shutil.which('soffice')
        if result:
            return result
        raise FileNotFoundError(
            "LibreOffice not found on Windows. "
            "Install from https://www.libreoffice.org/download/"
        )

    elif system == 'Linux':
        result = shutil.which('soffice') or shutil.which('libreoffice')
        if result:
            return result
        raise FileNotFoundError(
            "LibreOffice not found on Linux. "
            "Install with: sudo apt install libreoffice"
        )

    elif system == 'Darwin':
        candidates = [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
        ]
        for path in candidates:
            if os.path.exists(path):
                return path
        result = shutil.which('soffice')
        if result:
            return result
        raise FileNotFoundError(
            "LibreOffice not found on macOS. "
            "Install from https://www.libreoffice.org/download/"
        )

    raise FileNotFoundError(f"Unsupported OS: {system}")


LIBREOFFICE_PATH = None  # Lazily resolved


def _get_libreoffice():
    global LIBREOFFICE_PATH
    if LIBREOFFICE_PATH is None:
        LIBREOFFICE_PATH = _find_libreoffice()
        print(f"[LibreOffice] Found at: {LIBREOFFICE_PATH}")
    return LIBREOFFICE_PATH


# =====================================================
# TEXT REPLACEMENT IN PPTX
# =====================================================

def _replace_text_in_shape(shape, replacements: dict):
    """
    Replace placeholder text in a PowerPoint shape.

    replacements: dict of { "{{PLACEHOLDER}}": ("value", make_bold) }

    For each run containing a placeholder the run is SPLIT at the placeholder
    boundary into three sub-runs (all cloned from the original run's formatting):

        [before | value | after]

    The 'value' sub-run gets typeface "Poppins Bold" + b="1" when make_bold is
    True, or typeface "Poppins" with no b attribute when make_bold is False.
    The 'before' and 'after' sub-runs always get typeface "Poppins" with no b
    attribute, so surrounding sentence text stays non-bold regardless of how
    the original paragraph was formatted.
    """
    if not shape.has_text_frame:
        return

    import copy
    from lxml import etree

    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    QNAME_R   = f'{{{NS}}}r'
    QNAME_RPR = f'{{{NS}}}rPr'
    QNAME_T   = f'{{{NS}}}t'
    QNAME_LAT = f'{{{NS}}}latin'

    def _make_sub_run(source_r, text: str, bold: bool):
        """Clone source_r, set its text, and apply bold/normal typeface."""
        r = copy.deepcopy(source_r)
        # Set text
        t = r.find(QNAME_T)
        if t is None:
            t = etree.SubElement(r, QNAME_T)
        t.text = text
        # Adjust rPr
        rPr = r.find(QNAME_RPR)
        if rPr is None:
            rPr = etree.SubElement(r, QNAME_RPR)
            r.insert(0, rPr)
        if bold:
            rPr.set('b', '1')
            typeface = 'Poppins Bold'
        else:
            rPr.attrib.pop('b', None)      # remove bold attribute entirely
            typeface = 'Poppins'
        # Update latin typeface
        lat = rPr.find(QNAME_LAT)
        if lat is None:
            lat = etree.SubElement(rPr, QNAME_LAT)
        lat.set('typeface', typeface)
        return r

    def _replace_in_para(para_elem, placeholder: str, value: str, make_bold: bool):
        """Replace ONE placeholder in all runs of a paragraph, splitting as needed."""
        # Iterate over current <a:r> elements; rebuild list after each mutation
        processed = 0
        while True:
            runs = para_elem.findall(QNAME_R)
            found = False
            for r_elem in runs[processed:]:
                t_elem = r_elem.find(QNAME_T)
                text = (t_elem.text or '') if t_elem is not None else ''
                if placeholder not in text:
                    processed += 1
                    continue
                # Found — split this run
                idx   = text.index(placeholder)
                before = text[:idx]
                after  = text[idx + len(placeholder):]
                parent = r_elem.getparent()
                pos    = list(parent).index(r_elem)
                parent.remove(r_elem)
                insert_at = pos
                if before:
                    parent.insert(insert_at, _make_sub_run(r_elem, before, False))
                    insert_at += 1
                parent.insert(insert_at, _make_sub_run(r_elem, value, make_bold))
                insert_at += 1
                if after:
                    parent.insert(insert_at, _make_sub_run(r_elem, after, False))
                found = True
                break   # restart scan from the same position
            if not found:
                break

    txBody = shape.text_frame._txBody
    NS_P = f'{{{NS}}}p'
    for para_elem in txBody.findall(NS_P):
        for placeholder, (new_value, make_bold) in replacements.items():
            _replace_in_para(para_elem, placeholder, str(new_value), make_bold)


# =====================================================
# FILL PPTX TEMPLATE
# =====================================================

def fill_template(template_filename: str, replacements: dict, output_name: str) -> str:
    """
    Fill a PPTX template with placeholder values.
    
    Args:
        template_filename: Name of the template file in templates/ dir
        replacements: dict of { "{{PLACEHOLDER}}": ("value", make_bold) }
        output_name: Base name for the output file (no extension)
    
    Returns:
        Path to the filled PPTX file.
    """
    template_path = os.path.join(TEMPLATES_DIR, template_filename)
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_text_in_shape(shape, replacements)

    # Sanitize filename
    safe_name = (
        output_name
        .replace("/", "-").replace("\\", "-").replace(":", "-")
        .replace("*", "").replace("?", "").replace('"', "")
        .replace("<", "").replace(">", "").replace("|", "")
    )

    pptx_output = os.path.join(OUTPUT_DIR, f"{safe_name}.pptx")

    # Remove old file if exists
    if os.path.exists(pptx_output):
        os.remove(pptx_output)

    prs.save(pptx_output)
    print(f"[PPTX] Saved: {pptx_output}")

    return pptx_output


# =====================================================
# CONVERT PPTX TO PDF (LibreOffice Headless)
# =====================================================

def convert_to_pdf(pptx_path: str) -> str:
    """
    Convert a PPTX file to PDF using LibreOffice headless mode.
    Falls back to Microsoft PowerPoint via win32com on Windows if LibreOffice is not found.
    
    Args:
        pptx_path: Absolute path to the PPTX file.
    
    Returns:
        Path to the generated PDF file.
    """
    pdf_path = os.path.splitext(pptx_path)[0] + '.pdf'

    try:
        soffice = _get_libreoffice()
        output_dir = os.path.dirname(pptx_path)

        cmd = [
            soffice,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', output_dir,
            pptx_path
        ]

        print(f"[PDF] Converting via LibreOffice: {os.path.basename(pptx_path)}")

        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=60
        )

        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed.\n"
                f"stdout: {result.stdout}\n"
                f"stderr: {result.stderr}"
            )

        # Sometimes LibreOffice takes a moment to write the file
        for _ in range(10):
            if os.path.exists(pdf_path):
                break
            time.sleep(0.5)

        if not os.path.exists(pdf_path):
            raise FileNotFoundError(
                f"PDF was not generated. Expected at: {pdf_path}\n"
                f"LibreOffice output: {result.stdout}"
            )

        print(f"[PDF] Generated: {pdf_path}")
        return pdf_path

    except (FileNotFoundError, RuntimeError) as err:
        if platform.system() == 'Windows':
            print(f"[PDF] LibreOffice not available ({str(err)}). Trying win32com PowerPoint fallback...")
            try:
                import win32com.client
                import pythoncom
                
                # Initialize COM for the current thread (necessary in multi-threaded FastAPI)
                pythoncom.CoInitialize()
                
                # Use PowerPoint to convert
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                
                # Open presentation
                presentation = powerpoint.Presentations.Open(
                    os.path.abspath(pptx_path),
                    WithWindow=False
                )
                
                # Save as PDF (Format type 32 is PDF)
                presentation.SaveAs(os.path.abspath(pdf_path), 32)
                presentation.Close()
                powerpoint.Quit()
                
                # Clean up COM
                pythoncom.CoUninitialize()
                
                print(f"[PDF] Generated via win32com PowerPoint: {pdf_path}")
                return pdf_path
            except ImportError:
                raise FileNotFoundError(
                    "PDF conversion failed: LibreOffice was not found, and the win32com fallback is not available. "
                    "To fix this, please either install LibreOffice (https://www.libreoffice.org/download/) "
                    "or install the pywin32 library using: pip install pywin32"
                )
            except Exception as com_err:
                raise RuntimeError(
                    f"PDF conversion failed: LibreOffice is missing, and the Microsoft PowerPoint fallback failed. "
                    f"PowerPoint error: {str(com_err)}"
                )
        else:
            raise err


# =====================================================
# SEND EMAIL VIA GMAIL API
# =====================================================

def send_email(
    to_email: str,
    subject: str,
    body: str,
    attachment_path: str,
    cc_email: str | None = None
) -> None:
    """
    Send an email with a PDF attachment via Gmail API.
    
    Args:
        to_email: Recipient email address.
        subject: Email subject line.
        body: Plain text email body.
        attachment_path: Path to the PDF file to attach.
        cc_email: Optional CC email address.
    """
    service = get_gmail_service()

    message = EmailMessage()
    message['To'] = str(to_email).strip()
    message['From'] = SENDER_EMAIL
    if cc_email:
        message['Cc'] = cc_email
    message['Subject'] = subject
    message.set_content(body)

    # Attach file
    mime_type, _ = mimetypes.guess_type(attachment_path)
    if mime_type is None:
        mime_type = "application/pdf"
    main_type, sub_type = mime_type.split('/')

    with open(attachment_path, 'rb') as f:
        message.add_attachment(
            f.read(),
            maintype=main_type,
            subtype=sub_type,
            filename=os.path.basename(attachment_path)
        )

    raw_message = base64.urlsafe_b64encode(message.as_bytes()).decode()
    send_body = {'raw': raw_message}

    service.users().messages().send(userId="me", body=send_body).execute()
    print(f"[Email] Sent to: {to_email}")


# =====================================================
# GENERATE OFFER LETTER (Full Pipeline)
# =====================================================

def generate_offer_letter(
    employee_name: str,
    employee_email: str,
    employee_id: str,
    role_title: str,
    date: str,
    start_date: str,
    end_date: str,
    salary: str,
    responsibilities: str,
    template_filename: str = "OFFER_LETTER.pptx"
) -> dict:
    """
    Full pipeline: Fill template → Convert to PDF → Return paths.
    
    Returns:
        dict with 'pptx_path' and 'pdf_path' keys.
    """
    replacements = {
        "{{DATE}}": (date, True),
        "{{ID}}": (employee_id, True),
        "{{NAME}}": (employee_name, True),
        "{{ROLE}}": (role_title, True),
        "{{START_DATE}}": (start_date, False),
        "{{END_DATE}}": (end_date, False),
        "{{SALARY}}": (salary, False),
        "{{ROLES}}": (responsibilities, False),
    }

    output_name = f"{employee_name}_Offer_Letter"

    # Step 1: Fill PPTX template
    pptx_path = fill_template(template_filename, replacements, output_name)

    # Step 2: Convert to PDF
    pdf_path = convert_to_pdf(pptx_path)

    return {
        "pptx_path": pptx_path,
        "pdf_path": pdf_path,
    }


def send_offer_letter_email(
    employee_name: str,
    employee_email: str,
    salary: str,
    start_date: str,
    end_date: str,
    pdf_path: str,
    email_subject: str | None = None,
    email_body: str | None = None,
) -> None:
    """
    Send the offer letter PDF via email.
    
    If email_subject/email_body are not provided, uses defaults.
    """
    subject = email_subject or f"Offer Letter - {employee_name}"
    body = email_body or f"""
Dear {employee_name},

Congratulations!

We are delighted to offer you an internship position at Concept Simplified.

Please find your detailed Offer Letter attached to this email.

Acceptance Confirmation:
Reply to this email with "I accept the offer".

Important Notes:

• This internship will be conducted remotely, with 6 working days per week.

• The monthly stipend is ₹{salary}.

• Internship Duration:
  {start_date} to {end_date}

We are confident this opportunity will be a valuable step in your career journey.

Should you have any questions, feel free to contact us.

Warm regards,

Samkit Shah
Founder, Concept Simplified
"""

    send_email(
        to_email=employee_email,
        subject=subject,
        body=body,
        attachment_path=pdf_path,
        cc_email=CC_EMAIL
    )


# =====================================================
# GENERATE EXPERIENCE LETTER (Full Pipeline)
# =====================================================

def generate_experience_letter(
    employee_name: str,
    employee_email: str,
    employee_id: str,
    role: str,
    joining_date: str,
    relieving_date: str,
    duration: str,
    date: str,
    template_filename: str = "EXPERIENCE_LETTER.pptx"
) -> dict:
    """
    Full pipeline: Fill EXPERIENCE_LETTER template -> Convert to PDF -> Return paths.

    Returns:
        dict with 'pptx_path' and 'pdf_path' keys.
    """
    # Exact placeholders from the user's edited EXPERIENCE_LETTER.pptx:
    # TextBox 29: {{DATE}}
    # TextBox 30: {{ID}}
    # TextBox 32: {{NAME}}, {{ID}}, {{ROLE}}, {{JOIN_DATE}}, {{RELIEVE_DATE}}, {{DURATION}}
    replacements = {
        "{{DATE}}":         (date, True),
        "{{ID}}":           (employee_id, True),
        "{{NAME}}":         (employee_name, True),
        "{{ROLE}}":         (role, True),
        "{{JOIN_DATE}}":    (joining_date, True),
        "{{RELIEVE_DATE}}": (relieving_date, True),
        "{{DURATION}}":     (duration, True),
    }

    output_name = f"{employee_name}_Experience_Letter"

    pptx_path = fill_template(template_filename, replacements, output_name)
    pdf_path = convert_to_pdf(pptx_path)

    return {
        "pptx_path": pptx_path,
        "pdf_path": pdf_path,
    }


def send_experience_letter_email(
    employee_name: str,
    employee_email: str,
    role: str,
    joining_date: str,
    relieving_date: str,
    duration: str,
    pdf_path: str,
) -> None:
    """Send the experience letter PDF to the employee via Gmail."""
    subject = f"Experience Letter - {employee_name}"
    body = f"""Dear {employee_name},

Please find your Experience Letter attached to this email.

Details:
  Role: {role}
  Joining Date: {joining_date}
  Last Working Day: {relieving_date}
  Duration: {duration}

We sincerely appreciate your contributions to Concept Simplified and wish you all
the very best in your future endeavours.

Warm regards,

Samkit Shah
Founder, Concept Simplified
"""

    send_email(
        to_email=employee_email,
        subject=subject,
        body=body,
        attachment_path=pdf_path,
        cc_email=CC_EMAIL,
    )


# =====================================================
# GENERATE LOR (Full Pipeline)
# =====================================================

def generate_lor(
    employee_name: str,
    employee_email: str,
    employee_id: str,
    role: str,
    date: str,
    template_filename: str = "LOR.pptx"
) -> dict:
    """
    Full pipeline: Fill LOR template -> Convert to PDF -> Return paths.

    The recommendation body is a standard professional paragraph — no user
    input required. Only the dynamic fields (name, role) are substituted.

    Returns:
        dict with 'pptx_path' and 'pdf_path' keys.
    """
    # Exact placeholders from the user's edited LOR.pptx:
    # TextBox 29: {{DATE}}
    # TextBox 30: {{ID}}
    # TextBox 32: {{NAME}}, {{ROLE}}, {{RECOMMENDATION}}
    #
    # {{RECOMMENDATION}} receives a common professional paragraph — the HR
    # user does not need to write it manually.
    common_recommendation = (
        f"{employee_name} was a valued member of our team who demonstrated consistent "
        f"dedication, strong work ethic, and a collaborative attitude throughout their "
        f"tenure at Concept Simplified. Their contributions were meaningful and impactful, "
        f"and they approached every responsibility with professionalism and integrity. "
        f"We have observed their professional and personal growth over this period and are "
        f"fully confident in their capabilities. We believe {employee_name} will be a "
        f"valuable asset to any organisation or academic programme they choose to pursue."
    )

    replacements = {
        "{{DATE}}":           (date, True),
        "{{ID}}":             (employee_id, True),
        "{{NAME}}":           (employee_name, True),
        "{{ROLE}}":           (role, True),
        "{{RECOMMENDATION}}": (common_recommendation, False),
    }

    output_name = f"{employee_name}_LOR"

    pptx_path = fill_template(template_filename, replacements, output_name)
    pdf_path = convert_to_pdf(pptx_path)

    return {
        "pptx_path": pptx_path,
        "pdf_path": pdf_path,
    }


def send_lor_email(
    employee_name: str,
    employee_email: str,
    pdf_path: str,
) -> None:
    """Send the LOR PDF to the employee via Gmail."""
    subject = f"Letter of Recommendation - {employee_name}"
    body = f"""Dear {employee_name},

Please find your Letter of Recommendation attached to this email.

We are happy to support your journey and wish you the best of luck!

Warm regards,

Samkit Shah
Founder, Concept Simplified
"""

    send_email(
        to_email=employee_email,
        subject=subject,
        body=body,
        attachment_path=pdf_path,
        cc_email=CC_EMAIL,
    )
