"""
create_templates.py
===================
Creates EXPERIENCE_LETTER.pptx and LOR.pptx from the existing OFFER_LETTER.pptx.

Uses the first slide of OFFER_LETTER.pptx as a base (same letterhead, logo,
colours, footer, signature, stamp). Replaces only:
  - TextBox 31  →  letter heading ("INTERNSHIP LETTER" → new title)
  - TextBox 32  →  main body text with new placeholders

Run once:
    python create_templates.py
"""

import os
import copy
import shutil
from pptx import Presentation
from pptx.util import Pt
from lxml import etree

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATES_DIR = os.path.join(BASE_DIR, 'templates')
SOURCE = os.path.join(TEMPLATES_DIR, 'OFFER_LETTER.pptx')


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_run_text(run, text: str):
    """Set run text without disturbing XML namespace."""
    run.text = text


def _clear_textframe_and_set(tf, paragraphs_data):
    """
    Replace all paragraphs in a text frame.

    paragraphs_data is a list of lists of (text, bold) tuples per paragraph.
    Empty list entry → blank paragraph.
    """
    # Remove all existing paragraphs except the first (we need at least one)
    txBody = tf._txBody
    existing_paras = txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p')

    # Keep first para as template for formatting, remove rest
    first_para_xml = copy.deepcopy(existing_paras[0]) if existing_paras else None
    for p in existing_paras[1:]:
        txBody.remove(p)

    # Clear content of first paragraph
    if existing_paras:
        for r in existing_paras[0].findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
            existing_paras[0].remove(r)

    # Build new paragraphs
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    def make_paragraph(runs):
        """runs: list of (text, bold)"""
        p = copy.deepcopy(first_para_xml) if first_para_xml is not None else etree.SubElement(txBody, f'{{{ns}}}p')
        # Clear existing runs
        for r in p.findall(f'{{{ns}}}r'):
            p.remove(r)
        for text, bold in runs:
            r_elem = etree.SubElement(p, f'{{{ns}}}r')
            rPr = etree.SubElement(r_elem, f'{{{ns}}}rPr', attrib={'lang': 'en-US', 'dirty': '0'})
            if bold:
                rPr.set('b', '1')
            # Font
            solidFill = etree.SubElement(rPr, f'{{{ns}}}solidFill')
            srgbClr = etree.SubElement(solidFill, f'{{{ns}}}srgbClr', attrib={'val': '000000'})
            latin = etree.SubElement(rPr, f'{{{ns}}}latin', attrib={'typeface': 'Poppins Bold' if bold else 'Poppins', 'panose': '020B0502040204030204', 'pitchFamily': '34', 'charset': '0'})
            t_elem = etree.SubElement(r_elem, f'{{{ns}}}t')
            t_elem.text = text
        return p

    # Replace first para
    if paragraphs_data:
        first_runs = paragraphs_data[0]
        p0 = txBody.findall(f'{{{ns}}}p')[0]
        for r in p0.findall(f'{{{ns}}}r'):
            p0.remove(r)
        for text, bold in first_runs:
            r_elem = etree.SubElement(p0, f'{{{ns}}}r')
            rPr = etree.SubElement(r_elem, f'{{{ns}}}rPr', attrib={'lang': 'en-US', 'dirty': '0'})
            if bold:
                rPr.set('b', '1')
            solidFill = etree.SubElement(rPr, f'{{{ns}}}solidFill')
            srgbClr = etree.SubElement(solidFill, f'{{{ns}}}srgbClr', attrib={'val': '000000'})
            latin = etree.SubElement(rPr, f'{{{ns}}}latin', attrib={'typeface': 'Poppins Bold' if bold else 'Poppins'})
            t_elem = etree.SubElement(r_elem, f'{{{ns}}}t')
            t_elem.text = text

        # Add remaining paragraphs
        for runs in paragraphs_data[1:]:
            if not runs:
                p = make_paragraph([('', False)])
            else:
                p = make_paragraph(runs)
            txBody.append(p)


def _replace_all_text_in_shape(shape, old_text: str, new_text: str):
    """Simple full-text replacement across all runs in a text frame."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)


def _set_textbox_single_line(shape, text: str, bold: bool = True):
    """Replace all text in a textbox with a single run."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    # Set first run of first paragraph
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        tf.paragraphs[0].runs[0].font.bold = bold


def build_experience_letter_template():
    """Duplicate OFFER_LETTER slide 1, change title + body for Experience Letter."""
    prs = Presentation(SOURCE)

    # Remove slide 2 (index 1) — keep only the letterhead slide
    xml_slides = prs.slides._sldIdLst
    # We need to remove the second slide properly
    slide_to_remove = prs.slides[1]
    rId = prs.slides._sldIdLst[1].get('r:id') if len(prs.slides) > 1 else None
    if rId:
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[1]

    slide = prs.slides[0]

    # Find key textboxes
    tb_title = None  # TextBox 31 - "INTERNSHIP LETTER"
    tb_body = None   # TextBox 32 - main body
    tb_date = None   # TextBox 29 - DATE
    tb_id = None     # TextBox 30 - ID

    for shape in slide.shapes:
        if shape.name == 'TextBox 31':
            tb_title = shape
        elif shape.name == 'TextBox 32':
            tb_body = shape
        elif shape.name == 'TextBox 29':
            tb_date = shape
        elif shape.name == 'TextBox 30':
            tb_id = shape

    # Update title
    if tb_title:
        _set_textbox_single_line(tb_title, 'EXPERIENCE LETTER', bold=True)

    # Update date textbox (keep same placeholder format)
    if tb_date:
        for para in tb_date.text_frame.paragraphs:
            for run in para.runs:
                if '{{DATE}}' in run.text:
                    pass  # Keep as-is

    # Update ID textbox
    if tb_id:
        for para in tb_id.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace('ID : {{ID}}', 'ID : {{ID}}')  # Keep as-is

    # Replace body text — clear existing and write new experience letter body
    if tb_body:
        # Clear all runs and paragraphs and write fresh content
        tf = tb_body.text_frame
        # Build new content by modifying existing runs
        paragraphs_content = [
            "Dear {{NAME}},\n\nThis is to certify that {{NAME}} (Employee ID: {{ID}}) has worked with Concept Simplified as {{ROLE}} from {{JOIN_DATE}} to {{RELIEVE_DATE}}, a period of {{DURATION}}.\n\nDuring the tenure at Concept Simplified, {{NAME}} demonstrated exceptional professionalism, dedication, and a strong work ethic. They consistently delivered high-quality results and made meaningful contributions to their role and to the broader team.\n\nWe found {{NAME}} to be sincere, responsible, and a collaborative team player who approached every task with enthusiasm and integrity.\n\nWe wish {{NAME}} all the very best in their future endeavours and have no hesitation in recommending them for any opportunities they pursue.\n\nYours Sincerely,\n\nSamkit Shah\nFounder, Concept Simplified"
        ]

        # Use the simpler approach: replace text run by run
        all_paras = tf.paragraphs
        # Build new text content replacing all runs
        new_body = (
            "Dear {{NAME}},\n\n"
            "This is to certify that {{NAME}} (Employee ID: {{ID}}) has worked with "
            "Concept Simplified as {{ROLE}} from {{JOIN_DATE}} to {{RELIEVE_DATE}}, "
            "a period of {{DURATION}}.\n\n"
            "During their tenure, {{NAME}} demonstrated exceptional professionalism, "
            "dedication, and a strong work ethic. They consistently delivered high-quality "
            "results and made meaningful contributions to their role and to the broader team.\n\n"
            "We found {{NAME}} to be sincere, responsible, and a collaborative team player "
            "who approached every task with enthusiasm and integrity.\n\n"
            "We wish {{NAME}} all the very best in their future endeavours and have no "
            "hesitation in recommending them for any opportunities they pursue.\n\n\n\n"
            "Yours Sincerely,\n\n"
            "Samkit Shah\n"
            "Founder, Concept Simplified"
        )

        # Clear and set via direct XML manipulation
        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        txBody = tf._txBody
        paras = txBody.findall(f'{ns}p')

        # Keep only first paragraph and clear its runs
        for p in paras[1:]:
            txBody.remove(p)

        first_p = txBody.findall(f'{ns}p')[0]
        for r in first_p.findall(f'{ns}r'):
            first_p.remove(r)

        # Create a single run with the full text (line breaks become \n in text, we need separate paragraphs)
        lines = new_body.split('\n')
        pPr_elem = first_p.find(f'{ns}pPr')

        def create_para_with_text(ref_para, text, bold=False):
            """Create a new <a:p> element with a run."""
            new_p = copy.deepcopy(ref_para)
            for r in new_p.findall(f'{ns}r'):
                new_p.remove(r)
            if text:
                r_elem = etree.SubElement(new_p, f'{ns}r')
                rPr = etree.SubElement(r_elem, f'{ns}rPr', attrib={'lang': 'en-US', 'dirty': '0', 'sz': '1200'})
                if bold:
                    rPr.set('b', '1')
                solidFill = etree.SubElement(rPr, f'{ns}solidFill')
                etree.SubElement(solidFill, f'{ns}srgbClr', attrib={'val': '000000'})
                etree.SubElement(rPr, f'{ns}latin', attrib={'typeface': 'Poppins Bold' if bold else 'Poppins'})
                t_elem = etree.SubElement(r_elem, f'{ns}t')
                t_elem.text = text
            return new_p

        # Set first paragraph (Dear {{NAME}},)
        p0_run = etree.SubElement(first_p, f'{ns}r')
        p0_rPr = etree.SubElement(p0_run, f'{ns}rPr', attrib={'lang': 'en-US', 'dirty': '0', 'sz': '1200', 'b': '1'})
        solidFill0 = etree.SubElement(p0_rPr, f'{ns}solidFill')
        etree.SubElement(solidFill0, f'{ns}srgbClr', attrib={'val': '000000'})
        etree.SubElement(p0_rPr, f'{ns}latin', attrib={'typeface': 'Poppins Bold'})
        t0 = etree.SubElement(p0_run, f'{ns}t')
        t0.text = lines[0]

        # Add remaining lines as separate paragraphs
        for line in lines[1:]:
            new_p = create_para_with_text(first_p, line, bold=False)
            txBody.append(new_p)

    out_path = os.path.join(TEMPLATES_DIR, 'EXPERIENCE_LETTER.pptx')
    prs.save(out_path)
    print(f'[OK] Saved: {out_path}')


def build_lor_template():
    """Duplicate OFFER_LETTER slide 1, change title + body for Letter of Recommendation."""
    prs = Presentation(SOURCE)

    # Remove slide 2
    if len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[1]

    slide = prs.slides[0]

    tb_title = None
    tb_body = None
    tb_date = None
    tb_id = None

    for shape in slide.shapes:
        if shape.name == 'TextBox 31':
            tb_title = shape
        elif shape.name == 'TextBox 32':
            tb_body = shape
        elif shape.name == 'TextBox 29':
            tb_date = shape
        elif shape.name == 'TextBox 30':
            tb_id = shape

    # Update title
    if tb_title:
        _set_textbox_single_line(tb_title, 'LETTER OF RECOMMENDATION', bold=True)

    # Hide ID box (replace with empty or "Ref:")
    if tb_id:
        for para in tb_id.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace('ID : {{ID}}', 'Ref: {{EMPLOYEE_ID}}')

    # Replace body text for LOR
    if tb_body:
        new_body = (
            "To Whomsoever It May Concern,\n\n"
            "I am writing to highly recommend {{NAME}}.\n\n"
            "{{NAME}} has been associated with Concept Simplified as {{ROLE}}. During this "
            "period, they consistently demonstrated exceptional skills, a growth mindset, "
            "and the ability to deliver outstanding results.\n\n"
            "{{RECOMMENDATION}}\n\n"
            "I wholeheartedly endorse {{NAME}} and am confident they "
            "will bring the same level of excellence and commitment to your organisation. "
            "Please feel free to reach out for any further information.\n\n\n\n"
            "Yours Sincerely,\n\n"
            "Samkit Shah\n"
            "Founder, Concept Simplified"
        )

        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        tf = tb_body.text_frame
        txBody = tf._txBody
        paras = txBody.findall(f'{ns}p')

        for p in paras[1:]:
            txBody.remove(p)

        first_p = txBody.findall(f'{ns}p')[0]
        for r in first_p.findall(f'{ns}r'):
            first_p.remove(r)

        lines = new_body.split('\n')

        def create_para_with_text(ref_para, text, bold=False):
            new_p = copy.deepcopy(ref_para)
            for r in new_p.findall(f'{ns}r'):
                new_p.remove(r)
            if text:
                r_elem = etree.SubElement(new_p, f'{ns}r')
                rPr = etree.SubElement(r_elem, f'{ns}rPr', attrib={'lang': 'en-US', 'dirty': '0', 'sz': '1200'})
                if bold:
                    rPr.set('b', '1')
                solidFill = etree.SubElement(rPr, f'{ns}solidFill')
                etree.SubElement(solidFill, f'{ns}srgbClr', attrib={'val': '000000'})
                etree.SubElement(rPr, f'{ns}latin', attrib={'typeface': 'Poppins Bold' if bold else 'Poppins'})
                t_elem = etree.SubElement(r_elem, f'{ns}t')
                t_elem.text = text
            return new_p

        # Set first paragraph
        p0_run = etree.SubElement(first_p, f'{ns}r')
        p0_rPr = etree.SubElement(p0_run, f'{ns}rPr', attrib={'lang': 'en-US', 'dirty': '0', 'sz': '1200'})
        solidFill0 = etree.SubElement(p0_rPr, f'{ns}solidFill')
        etree.SubElement(solidFill0, f'{ns}srgbClr', attrib={'val': '000000'})
        etree.SubElement(p0_rPr, f'{ns}latin', attrib={'typeface': 'Poppins'})
        t0 = etree.SubElement(p0_run, f'{ns}t')
        t0.text = lines[0]

        for line in lines[1:]:
            new_p = create_para_with_text(first_p, line, bold=False)
            txBody.append(new_p)

    out_path = os.path.join(TEMPLATES_DIR, 'LOR.pptx')
    prs.save(out_path)
    print(f'[OK] Saved: {out_path}')


if __name__ == '__main__':
    print('Creating EXPERIENCE_LETTER.pptx...')
    build_experience_letter_template()

    print('Creating LOR.pptx...')
    build_lor_template()

    print('\nDone! Both templates created in:', TEMPLATES_DIR)
